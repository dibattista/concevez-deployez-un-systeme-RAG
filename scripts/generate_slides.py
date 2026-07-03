from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    # Création d'une nouvelle présentation
    prs = Presentation()

    # --- Slide 1 : Titre ---
    slide_layout = prs.slide_layouts[0] # 0 = Titre
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Pipeline RAG : Nettoyage & Préparation des Données"
    subtitle.text = "Justification des choix de Feature Selection\n(De 78 à 9 colonnes)"

    # --- Slide 2 : La Problématique des Données Brutes ---
    slide_layout = prs.slide_layouts[1] # 1 = Titre + Contenu
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "1. Constats sur les Données Brutes"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Une extraction de l'API OpenAgenda retourne 78 colonnes."
    
    p = tf.add_paragraph()
    p.text = "Sparsité massive :"
    p.level = 1
    p2 = tf.add_paragraph()
    p2.text = "Plus de 90% de valeurs manquantes sur de nombreux champs ('public', 'cibles', etc.)."
    p2.level = 2

    p = tf.add_paragraph()
    p.text = "Bruit multilingue :"
    p.level = 1
    p2 = tf.add_paragraph()
    p2.text = "Les champs textuels (titre, description) contiennent des dictionnaires avec plusieurs langues, ce qui perturbe le modèle (LLM)."
    p2.level = 2

    # --- Slide 3 : La Séparation Signal / Bruit ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "2. Séparation du Signal et du Bruit Technique"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Pourquoi supprimer des colonnes bien remplies (ex: timezone, slug, originAgenda) ?"
    
    p = tf.add_paragraph()
    p.text = "Bruit Applicatif vs. Pertinence Sémantique :"
    p.level = 1
    p2 = tf.add_paragraph()
    p2.text = "Ces métadonnées sont utiles pour le web, mais n'ont aucun poids sémantique pour répondre à une question utilisateur."
    p2.level = 2

    p = tf.add_paragraph()
    p.text = "Optimisation FAISS (Embeddings) :"
    p.level = 1
    p2 = tf.add_paragraph()
    p2.text = "Conserver uniquement le \"Cœur\" de l'événement (Titre, Description, Localisation, Dates) empêche de diluer le sens mathématique du vecteur."
    p2.level = 2

    # --- Slide 4 : Architecture de la Donnée Finale ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "3. Le Schéma de Données Retenu (9 colonnes)"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "La donnée a été épurée pour maximiser les performances du RAG :"
    
    p = tf.add_paragraph()
    p.text = "Signal Sémantique (Indexé dans FAISS) :"
    p.level = 1
    p2 = tf.add_paragraph()
    p2.text = "title_fr, description_fr, keywords_fr, location, firstTiming, lastTiming"
    p2.level = 2

    p = tf.add_paragraph()
    p.text = "Métadonnées de Restitution (Pour l'interface utilisateur) :"
    p.level = 1
    p2 = tf.add_paragraph()
    p2.text = "uid, slug, image (Non indexés, mais retournés pour l'affichage UI)"
    p2.level = 2

    # --- Slide 5 : Bilan Quantitatif ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "4. Bilan des Opérations de Nettoyage"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Impact direct sur le jeu de données :"
    
    p = tf.add_paragraph()
    p.text = "Déduplication : Suppression de centaines d'événements en double via l'identifiant 'uid'."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Filtre Localisation : Suppression des événements sans adresse (indispensable pour une recherche géographique)."
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Traduction : Extraction exclusive du français pour optimiser la fenêtre de contexte du LLM."
    p.level = 1

    # Sauvegarde
    prs.save('presentation_nettoyage.pptx')

if __name__ == "__main__":
    create_presentation()
